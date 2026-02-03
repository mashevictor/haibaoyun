import os
import sys

def extract_pptx(file_path):
    try:
        from pptx import Presentation
        
        if not os.path.exists(file_path):
            print(f"File not found: {file_path}")
            return

        prs = Presentation(file_path)
        text_content = []

        print(f"Extracting from {file_path}...")
        
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            
            if slide_text:
                text_content.append(f"=== Slide {i+1} ===\n" + "\n".join(slide_text) + "\n")
        
        full_text = "\n".join(text_content)
        
        output_file = "pptx_content.txt"
        with open(output_file, "w", encoding="utf-8") as f:
            f.write(full_text)
            
        print(f"Extracted {len(prs.slides)} slides to {output_file}")
        
    except ImportError:
        print("python-pptx not installed. Installing...")
        os.system(f"{sys.executable} -m pip install python-pptx")
        # Retry once
        try:
             from pptx import Presentation
             extract_pptx(file_path)
        except:
            print("Failed to run python-pptx after install attempt")

if __name__ == "__main__":
    extract_pptx("嗨豹云-融资-视频版.pptx")
